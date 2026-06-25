"""将动态视频嵌入PPT，作为封面页的全屏背景"""
from pptx import Presentation
from pptx.util import Inches, Emu, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_width = prs.slide_width
slide_height = prs.slide_height

# 第1页：封面 - 嵌入动态视频作为背景
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# 先放静态背景图作为后备
slide.shapes.add_picture(
    "/workspace/江底微光_版本A_1.png",
    0, 0, slide_width, slide_height
)

# 嵌入视频（PPT支持mp4，叠加在背景上或替换背景）
video_path = "/workspace/江底微光_动态循环.mp4"
# pptx的add_movie在较新版本支持
try:
    # 视频放全屏
    slide.shapes.add_movie(
        video_path,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5),
        poster_frame_image=None,
        mime_type='video/mp4'
    )
    print("视频嵌入成功（全屏）")
except Exception as e:
    print(f"视频嵌入方式1失败: {e}")
    # 备选：放在中央
    try:
        slide.shapes.add_movie(
            video_path,
            Inches(1.5), Inches(0.5),
            Inches(10.333), Inches(6.5),
            poster_frame_image=None,
            mime_type='video/mp4'
        )
        print("视频嵌入成功（居中）")
    except Exception as e2:
        print(f"视频嵌入方式2也失败: {e2}")

# 第2页：引语
slide2 = prs.slides.add_slide(slide_layout)
slide2.shapes.add_picture("/workspace/江底微光_版本A_2.png", 0, 0, slide_width, slide_height)

# 第3页：核心
slide3 = prs.slides.add_slide(slide_layout)
slide3.shapes.add_picture("/workspace/江底微光_版本A_3.png", 0, 0, slide_width, slide_height)

# 第4页：结尾
slide4 = prs.slides.add_slide(slide_layout)
slide4.shapes.add_picture("/workspace/江底微光_版本A_4.png", 0, 0, slide_width, slide_height)

output = "/workspace/江底微光_动态版.pptx"
prs.save(output)
print(f"PPT已生成: {output}")
print(f"文件大小: {os.path.getsize(output) / 1024:.0f} KB")
