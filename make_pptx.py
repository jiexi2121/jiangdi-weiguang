from pptx import Presentation
from pptx.util import Inches, Emu
import io

# 创建16:9演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

images = [
    "/workspace/江底微光_版本A_1.png",
    "/workspace/江底微光_版本A_2.png",
    "/workspace/江底微光_版本A_3.png",
    "/workspace/江底微光_版本A_4.png",
]

slide_width = prs.slide_width
slide_height = prs.slide_height

for img_path in images:
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    # 全屏铺满背景图
    slide.shapes.add_picture(img_path, 0, 0, slide_width, slide_height)

output_path = "/workspace/江底微光_版本A.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
