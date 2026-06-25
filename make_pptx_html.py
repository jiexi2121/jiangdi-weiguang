from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import os
import base64

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 读取HTML内容
with open("/workspace/江底微光_版本A.html", "r", encoding="utf-8") as f:
    html_content = f.read()

# 用data URI嵌入HTML到PPT（作为OLE对象不太可行，改用超链接方式）

# 方案：创建4页PPT，每页嵌入4张背景图，然后在每页添加一个操作按钮/超链接指向HTML文件
# 更实用的方案：把HTML作为附件嵌入PPT

images = [
    "/workspace/江底微光_版本A_1.png",
    "/workspace/江底微光_版本A_2.png",
    "/workspace/江底微光_版本A_3.png",
    "/workspace/江底微光_版本A_4.png",
]

slide_width = prs.slide_width
slide_height = prs.slide_height

for i, img_path in enumerate(images):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    # 全屏背景图
    slide.shapes.add_picture(img_path, 0, 0, slide_width, slide_height)
    
    # 右下角添加提示文字
    left = Inches(8.5)
    top = Inches(6.8)
    width = Inches(4.5)
    height = Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 双击同目录下的 江底微光_版本A.html 查看动态效果"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x80, 0xB8, 0xE0)
    p.font.italic = True
    p.alignment = PP_ALIGN.RIGHT

output_path = "/workspace/江底微光_动态版.pptx"
prs.save(output_path)

# 把HTML也复制一份方便一起发送
import shutil
shutil.copy("/workspace/江底微光_版本A.html", "/workspace/江底微光_动态版.html")

print(f"PPT已生成: {output_path}")
print(f"HTML已复制: /workspace/江底微光_动态版.html")
